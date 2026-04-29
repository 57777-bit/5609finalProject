"""
build_pptx_v2.py — CSCI 5609 Final Presentation (Editorial Earth revision)
Palette: warm cream bg, terracotta + sage accents, Georgia + Arial
17 slides: added "Updates Since FP3" as Slide 5
Output: presentation/CSCI5609_Final_Presentation_v2.pptx
"""
from __future__ import annotations
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import qrcode

# ─── Palette ──────────────────────────────────────────────────────────────────
BG          = RGBColor(0xF4, 0xEF, 0xE6)   # warm cream — all content slides
DARK_BG     = RGBColor(0x2A, 0x20, 0x18)   # warm dark  — title + Q&A
PRIMARY     = RGBColor(0x1F, 0x1B, 0x16)   # warm near-black
MUTED       = RGBColor(0x6B, 0x62, 0x59)   # captions, footers
TERRA       = RGBColor(0xB5, 0x53, 0x3C)   # terracotta accent
TERRA_LIGHT = RGBColor(0xD6, 0x7A, 0x5C)   # terracotta on dark bg
SAGE        = RGBColor(0x5F, 0x7A, 0x5A)   # sage green — Nordic / positive
HAIRLINE_C  = RGBColor(0xD9, 0xD0, 0xC2)   # thin dividers only
CREAM       = RGBColor(0xF4, 0xEF, 0xE6)   # cream text on dark slides

# ─── Typography ───────────────────────────────────────────────────────────────
SERIF = "Georgia"
SANS  = "Arial"
MONO  = "Courier New"

# ─── Dimensions ───────────────────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.6)
INNER_W = W - 2 * MARGIN       # ~12.133"

ROOT        = Path(__file__).resolve().parent
SCREENSHOTS = ROOT / "assets" / "screenshots"
QR_DIR      = ROOT / "assets" / "qr"
OUTPUT      = ROOT / "CSCI5609_Final_Presentation_v2.pptx"
LIVE_URL    = "https://57777-bit.github.io/5609finalProject/"
REPO_URL    = "https://github.com/57777-bit/5609finalProject"
TOTAL       = 17


# ─── Low-level helpers ────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height, *, wrap=True):
    """Create a textbox with zero internal margins."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    return tf


def run(tf_or_para, text, *, font=SANS, size=16, bold=False, italic=False,
        color=PRIMARY, tracking=0, caps=False, align=None):
    """Add a run to the last paragraph of a text frame (or a paragraph directly)."""
    from pptx.text.text import _Paragraph as PPTXPara
    if hasattr(tf_or_para, 'paragraphs'):
        p = tf_or_para.paragraphs[-1]
    else:
        p = tf_or_para
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if caps:
        r.font.all_caps = True
    if tracking:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(tracking))
    return r


def new_para(tf, *, align=PP_ALIGN.LEFT, line_spacing=None, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if space_before:
        p.space_before = Pt(space_before)
    return p


def hairline(slide, y, x0=None, x1=None):
    """Draw a thin horizontal divider."""
    if x0 is None: x0 = MARGIN
    if x1 is None: x1 = W - MARGIN
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, x1 - x0, Emu(9525))
    shp.fill.solid()
    shp.fill.fore_color.rgb = HAIRLINE_C
    shp.line.fill.background()


def eyebrow(slide, text, y=Inches(0.42), x=None, width=None, color=TERRA, dark=False):
    if x is None: x = MARGIN
    if width is None: width = INNER_W
    frame = tb(slide, x, y, width, Inches(0.32))
    run(frame, text, font=SANS, size=13, bold=True, caps=True, tracking=80,
        color=color)
    return frame


def title_text(slide, text, y=Inches(0.82), width=None, size=44, color=PRIMARY, dark=False):
    if width is None: width = Inches(9.5)
    frame = tb(slide, MARGIN, y, width, Inches(1.4), wrap=True)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    run(frame, text, font=SERIF, size=size, bold=True, color=color)
    return frame


def footer(slide, page_num: int, total: int = TOTAL):
    # Left: project title
    frame_l = tb(slide, MARGIN, H - Inches(0.42), Inches(6), Inches(0.32))
    run(frame_l, "The Geography of Opportunity", font=SANS, size=11, color=MUTED)
    # Right: page indicator
    frame_r = tb(slide, W - MARGIN - Inches(1.2), H - Inches(0.42), Inches(1.2), Inches(0.32))
    run(frame_r, f"{page_num:02d} / {total:02d}", font=SANS, size=11, color=MUTED,
        align=PP_ALIGN.RIGHT)


def demo_link(slide):
    """'→ live demo' terracotta italic top-right."""
    frame = tb(slide, W - MARGIN - Inches(2.0), Inches(0.40), Inches(2.0), Inches(0.32))
    run(frame, "→ live demo", font=SANS, size=13, italic=True, color=TERRA,
        align=PP_ALIGN.RIGHT)


def screenshot_or_placeholder(slide, x, y, w, filename, label):
    img_path = SCREENSHOTS / filename
    if img_path.exists():
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        target_h = int(w / aspect)
        pic = slide.shapes.add_picture(str(img_path), x, y, width=w, height=target_h)
        pic.line.color.rgb = HAIRLINE_C
        pic.line.width = Pt(0.5)
        return target_h
    # Placeholder rectangle
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(4.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = HAIRLINE_C
    shp.line.fill.background()
    frame = tb(slide, x + Inches(0.2), y + Inches(2.0), w - Inches(0.4), Inches(0.6))
    run(frame, f"[ {label} ]", font=SANS, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    return int(Inches(4.5))


def takeaway_block(slide, body_text, x, y, w):
    """TAKEAWAY section: hairline / label / Georgia Italic body / hairline."""
    hairline(slide, y, x, x + w)
    label_frame = tb(slide, x, y + Inches(0.18), w, Inches(0.30))
    run(label_frame, "TAKEAWAY", font=SANS, size=13, bold=True, caps=True,
        tracking=120, color=TERRA)
    body_frame = tb(slide, x, y + Inches(0.52), w, Inches(1.1), wrap=True)
    body_frame.vertical_anchor = MSO_ANCHOR.TOP
    run(body_frame, body_text, font=SERIF, size=20, italic=True, color=PRIMARY)
    bottom = y + Inches(1.72)
    hairline(slide, bottom, x, x + w)
    return bottom


def speaker_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def make_qr(url, fname):
    out = QR_DIR / fname
    if not out.exists():
        qrcode.make(url).save(out)
    return out


# ─── Slide builders ──────────────────────────────────────────────────────────

def build_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    # ── Right panel: 3D mobility map (red = high opportunity) ────────────────
    img_path = SCREENSHOTS / "step7_3d_red.png"
    if img_path.exists():
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        # Pin to full slide height; let width fall naturally to edge
        img_h = H
        img_w = Emu(int(img_h * iw / ih))
        img_x = W - img_w
        slide.shapes.add_picture(str(img_path), img_x, Emu(0), img_w, img_h)

        # Soft gradient vignette between text and image (semi-transparent dark rect)
        veil = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, img_x - Inches(0.5), Emu(0), Inches(1.2), H
        )
        veil.fill.solid()
        veil.fill.fore_color.rgb = DARK_BG
        veil.fill.fore_color.theme_color  # access to set alpha via XML
        # Set 60% transparency so it softly blends
        from lxml import etree
        sp_tree = veil._element
        solidFill = sp_tree.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                alpha_el = etree.SubElement(
                    srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                )
                alpha_el.set('val', '60000')  # 60 000 / 100 000 = 60% opacity
        veil.line.fill.background()

    # ── Left panel: text ─────────────────────────────────────────────────────
    TEXT_W = Inches(6.5)

    eyebrow(slide, "CSCI 5609 · Final Project · Spring 2026",
            y=Inches(1.55), color=TERRA_LIGHT, width=TEXT_W)

    # Title — stacks naturally to 2 lines at this width; looks deliberate
    frame = tb(slide, MARGIN, Inches(2.05), TEXT_W, Inches(2.6), wrap=True)
    run(frame, "The Geography of Opportunity",
        font=SERIF, size=52, bold=True, color=CREAM)

    # Subtitle
    frame = tb(slide, MARGIN, Inches(4.45), TEXT_W, Inches(1.0), wrap=True)
    run(frame, "How birthplace shapes the chance of climbing — and how the U.S. compares globally.",
        font=SERIF, size=20, italic=True, color=HAIRLINE_C)

    # Authors
    frame = tb(slide, MARGIN, Inches(5.7), TEXT_W, Inches(0.45))
    run(frame, "Brandon Borzello  ·  Harris Li  ·  Qi Wu  ·  Yiqi Huang",
        font=SANS, size=14, color=HAIRLINE_C)

    # URL
    frame = tb(slide, MARGIN, Inches(6.85), Inches(6.0), Inches(0.38))
    run(frame, LIVE_URL, font=MONO, size=11, color=TERRA_LIGHT)

    speaker_notes(slide,
        "We built an interactive scrollytelling piece asking two linked questions: "
        "how unequal is opportunity inside the U.S., and how does the U.S. compare "
        "against its wealthy peers? The argument is a local-to-global pivot. "
        "The 3D map on the right shows the full topography of opportunity — "
        "red peaks are counties where children of poor parents reach the highest adult incomes.")


def build_question(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Our Framing")

    # Dominant question in Georgia Italic 28pt
    frame = tb(slide, MARGIN, Inches(1.0), Inches(12.0), Inches(2.0), wrap=True)
    run(frame, ("How much does where an American child grows up shape their adult "
                "economic outcomes — and how does the U.S. compare against its wealthy peers?"),
        font=SERIF, size=28, italic=True, color=PRIMARY)

    hairline(slide, Inches(3.25))

    # Two acts
    frame = tb(slide, MARGIN, Inches(3.42), Inches(3.5), Inches(0.38))
    run(frame, "Two Acts", font=SANS, size=16, bold=True, caps=True,
        tracking=80, color=TERRA)

    # Part I
    frame = tb(slide, MARGIN, Inches(3.92), Inches(5.6), Inches(0.4))
    run(frame, "Part I — The American Problem", font=SANS, size=18, bold=True, color=PRIMARY)
    frame = tb(slide, MARGIN, Inches(4.38), Inches(5.6), Inches(0.6), wrap=True)
    p = frame.paragraphs[0]
    run(p, "Within the U.S., mobility is sharply ", font=SANS, size=16, color=PRIMARY)
    run(p, "geographic", font=SANS, size=16, bold=True, color=TERRA)
    run(p, ".", font=SANS, size=16, color=PRIMARY)

    hairline(slide, Inches(5.1), MARGIN, MARGIN + Inches(5.6))

    # Part II
    frame = tb(slide, MARGIN, Inches(5.25), Inches(5.6), Inches(0.4))
    run(frame, "Part II — The Global Perspective", font=SANS, size=18, bold=True, color=PRIMARY)
    frame = tb(slide, MARGIN, Inches(5.71), Inches(5.6), Inches(0.6), wrap=True)
    p = frame.paragraphs[0]
    run(p, "The U.S. as a whole sits ", font=SANS, size=16, color=PRIMARY)
    run(p, "below", font=SANS, size=16, bold=True, color=TERRA)
    run(p, " most of its wealthy peers.", font=SANS, size=16, color=PRIMARY)

    footer(slide, page)
    speaker_notes(slide,
        "Local-to-global structure lets the audience internalize within-country spread "
        "before weighing the U.S. against peers. School funding (Step 3) is the bridge.")


def build_data_sources(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Data Sources")
    title_text(slide, "All Public, All Static — No Backend")

    rows = [
        ("Opportunity Atlas",      "Chetty et al.",       "County-level adult outcomes by parent income"),
        ("GDIM 2023",              "World Bank",           "Intergenerational immobility (IGE) by country"),
        ("SWIID v9.6",             "Frederick Solt",       "Pre-/post-tax inequality (Gini)"),
        ("OECD Educ. at a Glance", "OECD",                "School funding centralization by country"),
        ("US Census TopoJSON",     "Mike Bostock",         "County + state boundaries (10M vertices)"),
    ]
    col_x = [MARGIN, MARGIN + Inches(3.8), MARGIN + Inches(6.5)]
    col_w = [Inches(3.6),    Inches(2.5),    Inches(5.4)]

    top = Inches(2.0)
    hairline(slide, top)

    # Column headers
    for i, (hdr, w) in enumerate(zip(["Dataset", "Attribution", "Use"], col_w)):
        frame = tb(slide, col_x[i], top + Inches(0.12), w, Inches(0.35))
        run(frame, hdr, font=SANS, size=12, bold=True, caps=True, tracking=80, color=TERRA)

    for ri, (ds, attr, use) in enumerate(rows):
        row_y = top + Inches(0.55) + ri * Inches(0.85)
        frame = tb(slide, col_x[0], row_y, col_w[0], Inches(0.65), wrap=True)
        run(frame, ds, font=SERIF, size=16, bold=True, color=PRIMARY)
        frame = tb(slide, col_x[1], row_y, col_w[1], Inches(0.65), wrap=True)
        run(frame, attr, font=SANS, size=14, color=MUTED)
        frame = tb(slide, col_x[2], row_y, col_w[2], Inches(0.65), wrap=True)
        run(frame, use, font=SANS, size=14, color=PRIMARY)

    hairline(slide, top + Inches(0.55) + len(rows) * Inches(0.85) + Inches(0.05))
    footer(slide, page)
    speaker_notes(slide, "All loaded statically — no backend. Bridge: GDIM × SWIID inner-join on country name.")


def build_tech_stack(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Tech Stack")
    title_text(slide, "Static Site, No Backend")

    cols = [
        ("Framework",      ["SvelteKit", "Svelte 5 runes ($state, $bindable, $props, $effect)"]),
        ("Visualization",  ["D3 v7 (scales, axes, projections)", "topojson-client", "deck.gl (3D bonus — lazy-loaded)"]),
        ("Build & Deploy", ["Vite dev server", "GitHub Actions → GitHub Pages (auto-deploy on push)"]),
    ]
    col_w = Inches(3.78)
    gap   = Inches(0.4)
    top   = Inches(2.0)

    for ci, (heading, bullets) in enumerate(cols):
        cx = MARGIN + ci * (col_w + gap)
        frame = tb(slide, cx, top, col_w, Inches(0.38))
        run(frame, heading, font=SANS, size=16, bold=True, caps=True, tracking=80, color=PRIMARY)
        frame = tb(slide, cx, top + Inches(0.5), col_w, Inches(3.0), wrap=True)
        tf = frame
        for bi, b in enumerate(bullets):
            p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            p.line_spacing = 1.5
            run(p, "–  " + b, font=SANS, size=16, color=PRIMARY)

    # Bottom note
    frame = tb(slide, MARGIN, Inches(5.85), INNER_W, Inches(0.65), wrap=True)
    run(frame, ("Sticky right panel + scrollytelling left panel — the chart stays anchored as "
                "the prose advances, so each narrative beat resolves against a visible mark."),
        font=SERIF, size=18, italic=True, color=MUTED)

    footer(slide, page)
    speaker_notes(slide, "Deploy is fully automated via GitHub Actions on every push to main.")


def build_updates(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Since Last Milestone")
    title_text(slide, "What Changed After FP3")

    # Lede
    frame = tb(slide, MARGIN, Inches(1.78), INNER_W, Inches(0.65), wrap=True)
    run(frame, ("Three rounds of peer review pushed us toward simpler interaction, "
                "tighter narrative continuity, and a 3D bonus that earns its place."),
        font=SERIF, size=18, italic=True, color=MUTED)

    updates = [
        ("01", "Play button replaces scroll trigger on Step 0",
         "Reviewers reported layout reflow races on fast scrolls — the map would re-mount "
         "mid-animation. We swapped the scroll trigger for an explicit Play / Replay button. "
         "Determinism over cleverness."),
        ("02", "Morph animation between Step 4 and Step 6",
         "The Gatsby Curve and the Mobility League ranking now share a transition: scatter dots "
         "morph into sorted bars (~1.2s, rect-with-rx interpolation). Same 20 countries, two "
         "views, one continuous argument."),
        ("03", "3D bonus added at Step 7 (deck.gl)",
         "A late addition. We held off until the 2D narrative was settled, then added the "
         "extruded county map as a 'magnitude reveal' — height = mobility, color = mobility — "
         "so the headline argument never depended on 3D for accessibility."),
    ]

    top = Inches(2.65)
    row_h = Inches(1.42)
    for i, (num, title, body) in enumerate(updates):
        y = top + i * row_h
        # Number
        frame = tb(slide, MARGIN, y, Inches(0.9), Inches(0.75))
        run(frame, num, font=SERIF, size=28, bold=True, color=TERRA)
        # Title + body
        frame = tb(slide, MARGIN + Inches(1.05), y, INNER_W - Inches(1.05),
                   Inches(1.3), wrap=True)
        tf = frame
        p0 = tf.paragraphs[0]
        run(p0, title, font=SANS, size=17, bold=True, color=PRIMARY)
        p1 = tf.add_paragraph()
        p1.line_spacing = 1.5
        p1.space_before = Pt(4)
        run(p1, body, font=SANS, size=15, color=PRIMARY)

    footer(slide, page)
    speaker_notes(slide,
        "Highlight the Play button change first — it's the most visible UX change from FP3. "
        "The morph animation is the engineering centrepiece. The 3D was intentionally held back "
        "to ensure it didn't undermine the accessible 2D narrative.")


# ─── Step slides ──────────────────────────────────────────────────────────────

STEPS = [
    dict(
        num=0, title="Scrollymap · The American Problem at a Glance",
        kicker="Step 0  ·  Part I — American Problem",
        lede="County-level mobility map with button-triggered animation through three tiers.",
        bullets=[
            [("Stuck ", {"color": TERRA, "bold": True}),
             ("(red) → Treading water (amber) → ", {}),
             ("Climbing ", {"color": SAGE, "bold": True}),
             ("(blue)", {})],
            "Toggleable to state-bubble (Dorling) view",
            "Click a state → zoom to its counties",
        ],
        takeaway="Play button replaces scroll trigger — fixes layout reflow races on fast scrolls.",
        screenshot="step0_scrollymap.png",
        notes="The opening uses a Play button instead of scroll trigger because scroll triggers were producing layout reflow races on fast scrolls.",
    ),
    dict(
        num=1, title="ScatterPlot · Poor vs Rich Children",
        kicker="Step 1  ·  Part I — American Problem",
        lede="Adult-rank by county: children of poor vs. rich parents.",
        bullets=[
            "Diagonal = perfect-mobility line",
            "Color encodes the gap (teal → yellow → red)",
            [("Annotation: ", {"color": TERRA, "bold": True}),
             ("the three counties with the largest mobility gap are labeled by name.", {})],
        ],
        takeaway="Distance from the diagonal is the mobility gap.",
        screenshot="step1_scatter.png",
        notes="The reddest, most-distant points get named on the chart so the prose has a target.",
    ),
    dict(
        num=2, title="ChangeMap · 1978 → 1992 Birth Cohorts",
        kicker="Step 2  ·  Part I — American Problem",
        lede="Animated transition from 1978 baseline to 1992 outcome.",
        bullets=[
            "Counties reveal in order of largest absolute change",
            [("Annotations: ", {"color": TERRA, "bold": True}),
             ("SOUTH / MIDWEST region tags anchor regional prose claims to the map.", {})],
            "Eye is drawn first to where mobility moved most",
        ],
        takeaway="Reveal order is the magnitude of change, not alphabetical.",
        screenshot="step2_changemap.png",
        notes="The animation reveals counties from the largest changes inward.",
    ),
    dict(
        num=3, title="SchoolFunding · The Bridge to Global",
        kicker="Step 3  ·  Part I — American Problem",
        lede="Stacked bar: Central/Federal vs. Local/State education funding share.",
        bullets=[
            "U.S. vs. five wealthy peers",
            [("Annotation: ", {"color": TERRA, "bold": True}),
             ("big '91% local' headline centered in the U.S. red segment.", {})],
            "Bridge from the American problem to the global perspective",
        ],
        takeaway="The U.S. funds schools through local property tax — geography of opportunity is geography of the tax base.",
        screenshot="step3_schoolfunding.png",
        notes="This is the pivot slide from the American problem to the global perspective.",
    ),
    dict(
        num=4, title="The Great Gatsby Curve",
        kicker="Step 4  ·  Part II — Global Perspective",
        lede="Cross-country scatter: after-tax Gini vs. intergenerational immobility (IGE).",
        bullets=[
            "Trend line: more inequality ↔ less mobility",
            [("Country labels: ", {"color": TERRA, "bold": True}),
             ("U.S. in terracotta, NORDIC BENCHMARKS caption in sage.", {})],
            [("Takeaway callout: ", {"color": TERRA, "bold": True}),
             ("'U.S. — high inequality, low mobility' with dashed leader.", {})],
        ],
        takeaway="The U.S. sits up and to the right — the worst quadrant on both axes.",
        screenshot="step4_gatsby.png",
        notes="The canonical Great Gatsby curve. The U.S. sits in the worst quadrant on both axes.",
    ),
    dict(
        num=5, title="Redistribution Dumbbell",
        kicker="Step 5  ·  Part II — Global Perspective",
        lede="For 14 countries: market Gini → disposable Gini after taxes/transfers.",
        bullets=[
            "Each row's Δ is labeled in points",
            [("Nordic accent: ", {"color": SAGE, "bold": True}),
             ("sage color + thicker line for visual binding to prose claims.", {})],
            "Headline contrasts U.S. redistribution against the strongest Nordic redistributor",
        ],
        takeaway="Denmark cuts ~25 points; the U.S. cuts under 10. The institutional choice matters.",
        screenshot="step5_dumbbell.png",
        notes="Same starting market inequality can land at very different places.",
    ),
    dict(
        num=6, title="Mobility League · The Closing Frame",
        kicker="Step 6  ·  Part II — Global Perspective",
        lede="Ranking of 20 major economies by intergenerational immobility (IGE).",
        bullets=[
            [("Morph animation: ", {"color": TERRA, "bold": True}),
             ("scatter dots morph into sorted bar ranking (~1.2s, rect-with-rx interpolation).", {})],
            [("U.S. callout: ", {"color": TERRA, "bold": True}),
             ("'rank 10 of 20' with dashed leader line.", {})],
            "▸ peer tags on Canada, Sweden, U.K., Japan",
        ],
        takeaway="Continuity from Step 4 to Step 6 is the whole argument: same world, viewed two ways.",
        screenshot="step6_league.png",
        notes="The morph makes visible that the SAME 20 countries are now ranked. Continuity is the argument.",
    ),
    dict(
        num=7, title="Bonus · A Topography of Opportunity",
        kicker="Bonus Step 7  ·  deck.gl",
        lede="3D extruded county map — height = mobility for poor-parent children (1992 cohort).",
        bullets=[
            [("Tech: ", {"color": TERRA, "bold": True}),
             ("deck.gl OrbitView + ColumnLayer + PathLayer, lazy-loaded.", {})],
            "Drag to orbit  ·  Scroll to zoom  ·  Hover for county detail",
            "2D answers WHICH counties are stuck. 3D answers HOW DRAMATIC.",
        ],
        takeaway="2D answers which counties are stuck. 3D makes the magnitude visceral.",
        screenshot="step7_3d.png",
        notes="Kept as bonus rather than opener — color encoding is more accessible for the headline argument.",
    ),
]


def build_step(prs, step: dict, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    eyebrow(slide, step["kicker"])
    demo_link(slide)

    # Title (36pt for step slides)
    title_text(slide, step["title"], y=Inches(0.82), size=36, width=Inches(9.0))

    # Lede
    frame = tb(slide, MARGIN, Inches(1.72), Inches(7.0), Inches(0.6), wrap=True)
    run(frame, step["lede"], font=SERIF, size=17, italic=True, color=MUTED)

    # Bullets
    frame = tb(slide, MARGIN, Inches(2.45), Inches(7.0), Inches(2.6), wrap=True)
    tf = frame
    for bi, item in enumerate(step["bullets"]):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.line_spacing = 1.55
        if bi > 0:
            p.space_before = Pt(4)
        # Bullet dash
        dr = p.add_run()
        dr.text = "–  "
        dr.font.name = SANS
        dr.font.size = Pt(18)
        dr.font.color.rgb = MUTED
        if isinstance(item, list):
            for text, opts in item:
                r2 = p.add_run()
                r2.text = text
                r2.font.name = opts.get("font", SANS)
                r2.font.size = Pt(18)
                r2.font.bold = opts.get("bold", False)
                r2.font.italic = opts.get("italic", False)
                r2.font.color.rgb = opts.get("color", PRIMARY)
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.name = SANS
            r2.font.size = Pt(18)
            r2.font.color.rgb = PRIMARY

    # TAKEAWAY
    takeaway_block(slide, step["takeaway"], MARGIN, Inches(5.3), Inches(7.0))

    # Screenshot — right column (40% of usable width)
    img_x = MARGIN + Inches(7.4)
    img_w = Inches(5.2)
    screenshot_or_placeholder(slide, img_x, Inches(1.1), img_w,
                               step["screenshot"], f"Step {step['num']}")

    footer(slide, page)
    speaker_notes(slide, step["notes"])


# ─── Design Decisions ─────────────────────────────────────────────────────────

DECISIONS = [
    ("Two-act structure",
     "Local → global lets the reader internalize within-country spread before weighing the U.S. against peers."),
    ("Sticky right panel",
     "Chart stays anchored; prose advances on the left. Each narrative beat resolves against a visible mark."),
    ("Word–image binding",
     "When prose names a country, that row is visibly distinguished — deep-blue label, peer tag, leader line."),
    ("Curated league of 20",
     "Restricted to major economies. Low-IGE outliers (Mauritius, Maldives) muddied the rich-peer framing."),
    ("Button trigger over scroll trigger",
     "Step 0 uses a Play button — fixes layout reflow races on fast scrolls. Determinism over cleverness."),
    ("3D as bonus, not opener",
     "Color encoding is more accessible for the headline argument. 3D is for magnitude impact at the end."),
]


def build_design_decisions(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Design Decisions")
    title_text(slide, "What Mattered — and Why")

    col_w = Inches(5.75)
    gap   = Inches(0.6)
    col_x = [MARGIN, MARGIN + col_w + gap]
    row_h = Inches(1.62)
    top   = Inches(1.9)

    for i, (title, body) in enumerate(DECISIONS):
        col = i % 2
        row = i // 2
        x = col_x[col]
        y = top + row * row_h

        # Row hairline (except first row)
        if row > 0 and col == 0:
            hairline(slide, y - Inches(0.08), MARGIN, W - MARGIN)

        frame = tb(slide, x, y, col_w, Inches(0.5))
        run(frame, f"{i+1}.", font=SERIF, size=22, bold=True, color=TERRA)

        frame = tb(slide, x + Inches(0.55), y + Inches(0.02), col_w - Inches(0.55), Inches(0.42))
        run(frame, title, font=SANS, size=16, bold=True, color=PRIMARY)

        frame = tb(slide, x, y + Inches(0.52), col_w, Inches(1.0), wrap=True)
        tf = frame
        p = tf.paragraphs[0]
        p.line_spacing = 1.5
        run(p, body, font=SANS, size=15, color=PRIMARY)

    footer(slide, page)
    speaker_notes(slide,
        "Most of these decisions came from things that broke. The button trigger replaced a "
        "scroll trigger that raced. The curated league replaced an 'all countries' view that "
        "buried the U.S. in noise.")


# ─── Animation Highlights ────────────────────────────────────────────────────

HIGHLIGHTS = [
    ("Morph animation",        "scatter → bar ranking on Step 6 (~1.2s, rect-with-rx interpolation)"),
    ("Animated reveal",        "county-by-county fill transition on Step 2"),
    ("Button-triggered play",  "Step 0 with Replay — deterministic, no scroll races"),
    ("Hover tooltips",         "every chart — country / county / value detail"),
    ("State drill-down",       "Dorling bubble view → click state → zoom to counties"),
    ("3D orbit + zoom",        "bonus Step 7 (deck.gl OrbitView, lazy-loaded)"),
    ("In-chart annotations",   "every step — prose references resolve to specific marks"),
]


def build_animation_highlights(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Animation & Interaction")
    title_text(slide, "What Moves, What the Reader Can Do")

    top = Inches(1.95)
    row_h = Inches(0.67)
    num_w = Inches(0.55)
    name_w = Inches(4.8)
    desc_x = MARGIN + num_w + name_w + Inches(0.3)
    desc_w = W - MARGIN - desc_x

    for i, (name, desc) in enumerate(HIGHLIGHTS):
        y = top + i * row_h
        if i > 0:
            hairline(slide, y - Inches(0.03), MARGIN, W - MARGIN)

        frame = tb(slide, MARGIN, y + Inches(0.04), num_w, Inches(0.5))
        run(frame, str(i + 1), font=SERIF, size=18, color=MUTED)

        frame = tb(slide, MARGIN + num_w, y + Inches(0.04), name_w, Inches(0.5))
        run(frame, name, font=SANS, size=16, bold=True, color=PRIMARY)

        frame = tb(slide, desc_x, y + Inches(0.04), desc_w, Inches(0.5), wrap=True)
        run(frame, desc, font=SANS, size=16, color=PRIMARY)

    footer(slide, page)
    speaker_notes(slide,
        "Step 6's morph is the visual climax — the same 20 countries from Step 4's scatter "
        "literally flow into a sorted ranking. Step 7's 3D adds a magnitude dimension.")


# ─── Reflection ───────────────────────────────────────────────────────────────

REFLECTIONS = [
    ("What worked",    SAGE,   ("Scrollytelling kept the argument sequenced. Sticky charts gave each prose beat "
                                "a visual anchor. The 3D bonus at the end gave a magnitude reveal without "
                                "competing with the 2D choropleth's accessibility.")),
    ("What we'd redo", MUTED,  ("Invest more in keyboard-only interaction for the 3D scene — orbit/zoom "
                                "currently rely on mouse gestures. Add an elevation slider for non-mouse users.")),
    ("Hardest part",   TERRA,  ("Keeping the same 20-country roster visually consistent across Steps 4, 5, 6. "
                                "The morph animation is what makes the roster legible — the eye carries "
                                "identity from one chart to the next.")),
]


def build_reflection(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    eyebrow(slide, "Reflection")
    title_text(slide, "What Worked · What We'd Redo")

    col_w = Inches(3.65)
    gap   = Inches(0.4)
    top   = Inches(1.95)

    for ci, (label, color, body) in enumerate(REFLECTIONS):
        cx = MARGIN + ci * (col_w + gap)

        # Vertical hairline between columns (not before first)
        if ci > 0:
            vl_x = cx - gap / 2
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, vl_x, top,
                                          Emu(9525), Inches(4.2))
            shp.fill.solid()
            shp.fill.fore_color.rgb = HAIRLINE_C
            shp.line.fill.background()

        frame = tb(slide, cx, top, col_w, Inches(0.38))
        run(frame, label, font=SANS, size=14, bold=True, caps=True,
            tracking=80, color=color)

        frame = tb(slide, cx, top + Inches(0.55), col_w, Inches(3.6), wrap=True)
        tf = frame
        p = tf.paragraphs[0]
        p.line_spacing = 1.5
        run(p, body, font=SERIF, size=16, color=PRIMARY)

    footer(slide, page)
    speaker_notes(slide,
        "If we had another week we'd add keyboard navigation to the 3D scene. "
        "The hardest engineering problem was keeping the 20-country roster consistent.")


# ─── Closing ──────────────────────────────────────────────────────────────────

def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)

    eyebrow(slide, "Q & A", y=Inches(1.55), color=TERRA_LIGHT)

    frame = tb(slide, MARGIN, Inches(2.05), Inches(12.0), Inches(1.6))
    run(frame, "Thank you.", font=SERIF, size=64, bold=True, color=CREAM)

    frame = tb(slide, MARGIN, Inches(3.55), Inches(10.0), Inches(0.65))
    run(frame, "We're happy to take questions.", font=SERIF, size=22, italic=True,
        color=HAIRLINE_C)

    # QR codes
    qr_live = make_qr(LIVE_URL, "live.png")
    qr_repo = make_qr(REPO_URL, "repo.png")
    qr_size = Inches(1.6)
    qr_y    = Inches(4.2)

    slide.shapes.add_picture(str(qr_live), Inches(1.2), qr_y, qr_size, qr_size)
    frame = tb(slide, Inches(1.2), qr_y + qr_size + Inches(0.1), qr_size, Inches(0.35))
    run(frame, "LIVE DEMO", font=SANS, size=14, bold=True, caps=True,
        tracking=80, color=TERRA_LIGHT, align=PP_ALIGN.CENTER)
    frame = tb(slide, Inches(0.8), qr_y + qr_size + Inches(0.5), Inches(2.4), Inches(0.35))
    run(frame, "57777-bit.github.io/5609finalProject", font=MONO, size=10,
        color=HAIRLINE_C, align=PP_ALIGN.CENTER)

    slide.shapes.add_picture(str(qr_repo), Inches(8.3), qr_y, qr_size, qr_size)
    frame = tb(slide, Inches(8.3), qr_y + qr_size + Inches(0.1), qr_size, Inches(0.35))
    run(frame, "REPO", font=SANS, size=14, bold=True, caps=True,
        tracking=80, color=TERRA_LIGHT, align=PP_ALIGN.CENTER)
    frame = tb(slide, Inches(7.8), qr_y + qr_size + Inches(0.5), Inches(2.6), Inches(0.35))
    run(frame, "github.com/57777-bit/5609finalProject", font=MONO, size=10,
        color=HAIRLINE_C, align=PP_ALIGN.CENTER)

    # Acknowledgements — prominent, not 9pt
    ack_y = H - Inches(1.15)
    hairline(slide, ack_y, MARGIN, W - MARGIN, )
    frame_lbl = tb(slide, MARGIN, ack_y + Inches(0.12), Inches(3.0), Inches(0.3))
    run(frame_lbl, "Acknowledgements", font=SANS, size=13, bold=True, caps=True,
        tracking=80, color=TERRA_LIGHT)
    frame_body = tb(slide, MARGIN, ack_y + Inches(0.45), INNER_W, Inches(0.55), wrap=True)
    p = frame_body.paragraphs[0]
    p.line_spacing = 1.5
    run(p, ("Chetty et al. · Opportunity Atlas  ·  World Bank GDIM 2023  ·  Frederick Solt SWIID v9.6  ·  "
            "Mike Bostock TopoJSON Atlas  ·  D3  ·  Svelte / SvelteKit  ·  deck.gl  ·  "
            "OECD Education at a Glance  ·  Special thanks to Prof. Qianwen Wang and the CSCI 5609 TAs."),
        font=SERIF, size=14, color=CREAM)

    speaker_notes(slide,
        "Open the floor. Common questions: why 1992 specifically (earliest robust Opp Atlas cohort), "
        "why IGE not rank correlation (GDIM uses IGE, kept their measure), why 3D is bonus not opener.")


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_title(prs)               # 1 — dark bg
    build_question(prs,    2)      # 2
    build_data_sources(prs, 3)     # 3
    build_tech_stack(prs,  4)      # 4
    build_updates(prs,     5)      # 5 — NEW
    for i, step in enumerate(STEPS):
        build_step(prs, step, 6 + i)   # 6–13
    build_design_decisions(prs,  14)   # 14
    build_animation_highlights(prs, 15)# 15
    build_reflection(prs,        16)   # 16
    build_closing(prs)                 # 17 — dark bg

    prs.save(OUTPUT)
    total = len(prs.slides)
    size  = OUTPUT.stat().st_size // 1024
    print(f"Saved  {OUTPUT.name}  ({total} slides, {size} KB)")


if __name__ == "__main__":
    main()
