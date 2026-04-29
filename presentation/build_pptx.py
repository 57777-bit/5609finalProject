"""
Generate the CSCI 5609 final-project presentation deck (.pptx) from PRESENTATION_OUTLINE.md.

Re-runnable: edits to outline, talking points, or screenshots take effect on next run.
Output: presentation/CSCI5609_Final_Presentation.pptx
"""

from __future__ import annotations

from pathlib import Path
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import qrcode

# ─────────────────────────────────────────────────────────────────────────────
# Design System (matches the live D3 site palette)
# ─────────────────────────────────────────────────────────────────────────────

NAVY = RGBColor(0x2C, 0x3E, 0x50)
US_RED = RGBColor(0xC0, 0x39, 0x2B)
NORDIC_GREEN = RGBColor(0x27, 0xAE, 0x60)
DATA_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BG = RGBColor(0xFA, 0xFB, 0xFD)
SLATE = RGBColor(0x7B, 0x8A, 0x8B)
BORDER_GRAY = RGBColor(0xE3, 0xE8, 0xEE)
DARK_INK = RGBColor(0x1A, 0x23, 0x2E)
SOFT_GRAY = RGBColor(0xAA, 0xB1, 0xBB)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "CSCI5609_Final_Presentation.pptx"
SCREENSHOTS = ROOT / "assets" / "screenshots"
QR_DIR = ROOT / "assets" / "qr"

LIVE_URL = "https://57777-bit.github.io/5609finalProject/"
REPO_URL = "https://github.com/57777-bit/5609finalProject"


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_solid_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=DARK_INK, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *, size=16,
                color=DARK_INK, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        runs = item if isinstance(item, list) else [(item, {})]
        prefix_run = p.add_run()
        prefix_run.text = "•  "
        prefix_run.font.name = FONT_BODY
        prefix_run.font.size = Pt(size)
        prefix_run.font.color.rgb = SLATE
        for text, opts in runs:
            r = p.add_run()
            r.text = text
            r.font.name = opts.get("font", FONT_BODY)
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", False)
            r.font.color.rgb = opts.get("color", color)
    return tb


def add_rect(slide, left, top, width, height, color, *, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    add_solid_fill(shp, color)
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    return shp


def set_slide_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_speaker_notes(slide, notes_text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_screenshot_or_placeholder(slide, left, top, width, height, name, label):
    img = SCREENSHOTS / name
    if img.exists():
        pic = slide.shapes.add_picture(str(img), left, top, width=width, height=height)
        # subtle border
        pic.line.color.rgb = BORDER_GRAY
        pic.line.width = Pt(0.75)
        return pic
    # Placeholder
    add_rect(slide, left, top, width, height, LIGHT_BG, line_color=BORDER_GRAY)
    add_text(slide, left, top + height // 2 - Inches(0.4), width, Inches(0.3),
             f"[ Screenshot · {label} ]", size=14, color=SOFT_GRAY,
             align=PP_ALIGN.CENTER)
    add_text(slide, left, top + height // 2, width, Inches(0.3),
             f"drop {name} into presentation/assets/screenshots/", size=10,
             color=SOFT_GRAY, font=FONT_MONO, align=PP_ALIGN.CENTER)


def slide_header(slide, title, *, kicker=None, demo_tag=False):
    """Standard top header band with optional kicker + DEMO tag."""
    # Header strip
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.6), NAVY)
    add_text(slide, Inches(0.6), Inches(0.12), Inches(11), Inches(0.36),
             "The Geography of Opportunity   ·   CSCI 5609 Spring 2026",
             size=11, color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_BODY)

    if kicker:
        add_text(slide, Inches(0.6), Inches(0.85), Inches(8), Inches(0.35),
                 kicker, size=12, bold=True, color=US_RED, font=FONT_BODY)

    title_top = Inches(1.15) if kicker else Inches(0.85)
    add_text(slide, Inches(0.6), title_top, Inches(10), Inches(0.7),
             title, size=32, bold=True, color=NAVY, font=FONT_TITLE)

    if demo_tag:
        # DEMO badge top right
        badge = add_rect(slide, Inches(11.5), Inches(0.85), Inches(1.35),
                         Inches(0.45), US_RED)
        tb = slide.shapes.add_textbox(Inches(11.5), Inches(0.85),
                                      Inches(1.35), Inches(0.45))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "▶ LIVE DEMO"
        r.font.name = FONT_TITLE
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def slide_footer(slide, page_num, total):
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), LIGHT_BG)
    add_text(slide, Inches(0.6), SLIDE_H - Inches(0.32), Inches(8), Inches(0.25),
             "Brandon Borzello · Harris Li · Qi Wu · Yiqi Huang", size=10,
             color=SLATE)
    add_text(slide, Inches(11.5), SLIDE_H - Inches(0.32), Inches(1.3), Inches(0.25),
             f"{page_num} / {total}", size=10, color=SLATE, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def build_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NAVY)

    # Decorative red accent bar
    add_rect(slide, Inches(0.6), Inches(2.2), Inches(0.15), Inches(2.5), US_RED)

    add_text(slide, Inches(0.95), Inches(2.0), Inches(11.5), Inches(0.5),
             "CSCI 5609 · FINAL PROJECT · SPRING 2026", size=14, bold=True,
             color=RGBColor(0xFF, 0xC0, 0xB6), font=FONT_TITLE)

    add_text(slide, Inches(0.95), Inches(2.55), Inches(12), Inches(1.4),
             "The Geography of Opportunity", size=58, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_TITLE)

    add_text(slide, Inches(0.95), Inches(3.85), Inches(12), Inches(0.7),
             "How birthplace shapes the chance of climbing —",
             size=24, color=RGBColor(0xE0, 0xE6, 0xEE), font=FONT_TITLE)
    add_text(slide, Inches(0.95), Inches(4.30), Inches(12), Inches(0.7),
             "and how the U.S. compares globally.",
             size=24, color=RGBColor(0xE0, 0xE6, 0xEE), font=FONT_TITLE)

    add_rect(slide, Inches(0.95), Inches(5.3), Inches(7), Inches(0.04),
             RGBColor(0xFF, 0xC0, 0xB6))

    add_text(slide, Inches(0.95), Inches(5.5), Inches(12), Inches(0.4),
             "Brandon Borzello   ·   Harris Li   ·   Qi Wu   ·   Yiqi Huang",
             size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_TITLE)

    add_text(slide, Inches(0.95), Inches(6.7), Inches(12), Inches(0.3),
             LIVE_URL, size=12, color=RGBColor(0xFF, 0xC0, 0xB6), font=FONT_MONO)

    set_speaker_notes(
        slide,
        "We built an interactive scrollytelling piece that asks two linked "
        "questions: how unequal is opportunity inside the U.S., and how does "
        "the U.S. as a whole stack up against its wealthy peers? The whole "
        "argument is a local-to-global pivot, and we kept it as a single "
        "scroll from a county map all the way out to a 20-country mobility "
        "ranking.",
    )


def build_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    slide_header(slide, "The Question", kicker="OUR FRAMING")

    # Center pull quote
    add_rect(slide, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0),
             RGBColor(0xFF, 0xFF, 0xFF), line_color=BORDER_GRAY)
    add_rect(slide, Inches(1.0), Inches(2.2), Inches(0.12), Inches(2.0), US_RED)

    add_text(slide, Inches(1.5), Inches(2.5), Inches(10.5), Inches(0.6),
             "How much does where an American child grows up shape their adult",
             size=22, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(slide, Inches(1.5), Inches(3.05), Inches(10.5), Inches(0.6),
             "economic outcomes — and how does the U.S. compare against",
             size=22, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(slide, Inches(1.5), Inches(3.6), Inches(10.5), Inches(0.6),
             "its wealthy peers?", size=22, bold=True, color=NAVY,
             font=FONT_TITLE)

    # Two acts
    add_text(slide, Inches(1.0), Inches(4.7), Inches(11), Inches(0.4),
             "Two acts", size=14, bold=True, color=US_RED, font=FONT_TITLE)

    add_rect(slide, Inches(1.0), Inches(5.2), Inches(5.6), Inches(1.6),
             RGBColor(0xFF, 0xFF, 0xFF), line_color=BORDER_GRAY)
    add_text(slide, Inches(1.3), Inches(5.4), Inches(5), Inches(0.4),
             "Part I — The American Problem", size=16, bold=True, color=NAVY)
    add_text(slide, Inches(1.3), Inches(5.85), Inches(5.2), Inches(0.9),
             "Within the U.S., mobility is sharply geographic.",
             size=14, color=SLATE)

    add_rect(slide, Inches(6.85), Inches(5.2), Inches(5.6), Inches(1.6),
             RGBColor(0xFF, 0xFF, 0xFF), line_color=BORDER_GRAY)
    add_text(slide, Inches(7.15), Inches(5.4), Inches(5.2), Inches(0.4),
             "Part II — The Global Perspective", size=16, bold=True, color=NAVY)
    add_text(slide, Inches(7.15), Inches(5.85), Inches(5.2), Inches(0.9),
             "The U.S. as a whole sits below most of its wealthy peers.",
             size=14, color=SLATE)

    set_speaker_notes(
        slide,
        "We deliberately structure as local→global so the audience first "
        "feels the within-country spread, then is asked to weigh the U.S. "
        "against peer economies. The pivot point is school funding — Step 3 "
        "is the bridge.",
    )


def build_data_sources(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    slide_header(slide, "Data Sources", kicker="ALL PUBLIC, ALL STATIC")

    sources = [
        ("Opportunity Atlas", "Chetty et al.", "County-level adult outcomes by parent income"),
        ("GDIM 2023", "World Bank", "Intergenerational immobility (IGE) by country"),
        ("SWIID v9.6", "Frederick Solt", "Pre-/post-tax inequality (Gini)"),
        ("OECD Education at a Glance", "OECD", "School funding centralization"),
        ("US Census TopoJSON Atlas", "Mike Bostock", "County + state boundaries (10M vertices)"),
    ]

    top = Inches(2.0)
    row_h = Inches(0.85)
    # Header row
    add_rect(slide, Inches(0.6), top, Inches(12.1), Inches(0.5), NAVY)
    add_text(slide, Inches(0.85), top + Inches(0.12), Inches(3.3), Inches(0.3),
             "DATASET", size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(4.5), top + Inches(0.12), Inches(2.8), Inches(0.3),
             "ATTRIBUTION", size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(7.5), top + Inches(0.12), Inches(5.0), Inches(0.3),
             "USE", size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    for i, (ds, attr, use) in enumerate(sources):
        y = top + Inches(0.5) + i * row_h
        bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF4, 0xF7, 0xFB)
        add_rect(slide, Inches(0.6), y, Inches(12.1), row_h, bg, line_color=BORDER_GRAY)
        add_text(slide, Inches(0.85), y + Inches(0.22), Inches(3.5), Inches(0.4),
                 ds, size=14, bold=True, color=NAVY)
        add_text(slide, Inches(4.5), y + Inches(0.22), Inches(2.8), Inches(0.4),
                 attr, size=12, color=SLATE)
        add_text(slide, Inches(7.5), y + Inches(0.22), Inches(5.0), Inches(0.4),
                 use, size=12, color=DARK_INK)

    set_speaker_notes(
        slide,
        "All public, all loaded statically — no backend. Counties at 10M-vertex "
        "Albers USA projection. We bridge GDIM and SWIID by inner-joining on "
        "country name; that's how Step 4 (Gatsby curve) and Step 5 "
        "(redistribution dumbbell) share the same 14-country roster.",
    )


def build_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    slide_header(slide, "Tech Stack", kicker="STATIC SITE, NO BACKEND")

    cards = [
        ("Framework", ["SvelteKit", "Svelte 5 runes ($state, $bindable, $props, $effect)"]),
        ("Visualization", ["D3 v7 (scales, axes, projections)", "topojson-client", "deck.gl (3D bonus)"]),
        ("Build & Deploy", ["Vite dev server", "GitHub Actions → GitHub Pages"]),
    ]

    card_w = Inches(3.95)
    card_h = Inches(3.4)
    gap = Inches(0.2)
    start_x = Inches(0.6)
    top = Inches(2.0)

    for i, (heading, items) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, top, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF),
                 line_color=BORDER_GRAY)
        add_rect(slide, x, top, card_w, Inches(0.08), US_RED)
        add_text(slide, x + Inches(0.3), top + Inches(0.3), card_w, Inches(0.4),
                 heading, size=16, bold=True, color=NAVY)
        add_bullets(slide, x + Inches(0.3), top + Inches(0.85),
                    card_w - Inches(0.4), card_h - Inches(1.0),
                    items, size=13, color=DARK_INK)

    add_text(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
             "Sticky right panel + scrollytelling left panel:",
             size=14, bold=True, color=NAVY)
    add_text(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(1.2),
             "the chart stays anchored as the prose advances, so each "
             "narrative beat resolves against a visible mark.", size=14,
             color=SLATE)

    set_speaker_notes(
        slide,
        "Static site — no backend. Sticky right panel with scrollytelling "
        "left panel. The chart stays anchored as the prose advances, so each "
        "narrative beat resolves against a visible mark. Deploy is fully "
        "automated via GitHub Actions on every push to main.",
    )


@dataclass
class StepSlide:
    step_num: int
    title: str
    summary: str
    bullets: list
    annotation: str | None
    talking_point: str
    screenshot: str
    bonus: bool = False


STEPS: list[StepSlide] = [
    StepSlide(
        step_num=0,
        title="Scrollymap · The American Problem at a Glance",
        summary="County-level mobility map with button-triggered animation through three tiers.",
        bullets=[
            [("Stuck (red) ", {"bold": True, "color": US_RED}),
             ("→ Treading water (amber) → ", {}),
             ("Climbing (blue)", {"bold": True, "color": DATA_BLUE})],
            "Toggleable to state-bubble (Dorling) view",
            "Click a state → zoom to its counties",
        ],
        annotation="Play button instead of scroll trigger — fixes layout reflow races.",
        talking_point=(
            "The opening uses a Play button instead of a scroll trigger because "
            "scroll triggers were producing layout reflow races on fast scrolls. "
            "The button locks the section to 100vh and removes that bug class."
        ),
        screenshot="step0_scrollymap.png",
    ),
    StepSlide(
        step_num=1,
        title="ScatterPlot · Poor vs Rich Children",
        summary="Adult-rank by county for poor-children vs rich-children parents.",
        bullets=[
            "Diagonal = perfect-mobility line",
            "Color encodes the gap (teal → yellow → red)",
            [("Annotation: ", {"bold": True, "color": US_RED}),
             ("the three counties with the largest mobility gap are labeled.", {})],
        ],
        annotation="Distance from the diagonal is the mobility gap.",
        talking_point=(
            "Distance from the diagonal is the mobility gap. The reddest, "
            "most-distant points get named on the chart so the prose has a target."
        ),
        screenshot="step1_scatter.png",
    ),
    StepSlide(
        step_num=2,
        title="ChangeMap · 1978 → 1992 Birth Cohorts",
        summary="Animated transition from the 1978 baseline to the 1992 outcome.",
        bullets=[
            "Counties reveal in order of largest absolute change",
            [("Annotations: ", {"bold": True, "color": US_RED}),
             ("SOUTH / MIDWEST region tags anchor regional prose claims.", {})],
            "Eye is drawn first to where mobility moved most",
        ],
        annotation="Reveal order is the magnitude of change, not alphabetical.",
        talking_point=(
            "The animation reveals counties from the largest changes inward, "
            "so the eye is drawn first to where mobility moved most. The "
            "regional tags let us say 'the Mississippi Delta lost ground' and "
            "have a target to point at."
        ),
        screenshot="step2_changemap.png",
    ),
    StepSlide(
        step_num=3,
        title="SchoolFunding · The Bridge to the Global View",
        summary="Stacked bar: Central/Federal vs. Local/State education funding share.",
        bullets=[
            "U.S. vs. five wealthy peers",
            [("Annotation: ", {"bold": True, "color": US_RED}),
             ("big 'X% local' headline centered in the U.S. red segment.", {})],
            "Pivot from the American problem to the global perspective",
        ],
        annotation="The U.S. funds schools through local property tax — geography of opportunity is geography of the tax base.",
        talking_point=(
            "This is the bridge from the American problem to the global "
            "perspective. The U.S. funds schools through local property tax — "
            "so the geography of opportunity is partly the geography of the "
            "tax base."
        ),
        screenshot="step3_schoolfunding.png",
    ),
    StepSlide(
        step_num=4,
        title="The Great Gatsby Curve",
        summary="Cross-country scatter: after-tax Gini vs. immobility (IGE).",
        bullets=[
            "Trend line: more inequality ↔ less mobility",
            [("Country labels: ", {"bold": True, "color": US_RED}),
             ("U.S. in red, NORDIC BENCHMARKS caption", {})],
            [("Takeaway callout: ", {"bold": True, "color": US_RED}),
             ("'U.S. — high inequality, low mobility' with dashed leader.", {})],
        ],
        annotation="The U.S. sits up and to the right — the worst quadrant on both axes.",
        talking_point=(
            "This is the canonical Great Gatsby curve. The U.S. sits up and "
            "to the right — the worst quadrant on both axes. The dashed leader "
            "line on the U.S. dot means the prose claim resolves to a specific "
            "mark."
        ),
        screenshot="step4_gatsby.png",
    ),
    StepSlide(
        step_num=5,
        title="Redistribution Dumbbell",
        summary="For 14 countries: market Gini → disposable Gini after taxes/transfers.",
        bullets=[
            "Each row's Δ is labeled in points",
            [("Nordic accent: ", {"bold": True, "color": NORDIC_GREEN}),
             ("green color + thicker line for visual binding to prose", {})],
            "Headline contrasts U.S. redistribution against the strongest Nordic redistributor",
        ],
        annotation="Denmark cuts ~25 points; the U.S. cuts under 10. The institutional choice matters.",
        talking_point=(
            "Same starting market inequality can land at very different "
            "places. Denmark cuts roughly 25 points; the U.S. cuts under 10. "
            "The institutional choice matters more than the starting "
            "condition."
        ),
        screenshot="step5_dumbbell.png",
    ),
    StepSlide(
        step_num=6,
        title="Mobility League · The Closing Frame",
        summary="Ranking of 20 major economies by intergenerational immobility (IGE).",
        bullets=[
            [("Animation: ", {"bold": True, "color": US_RED}),
             ("scatter dots morph into sorted bar ranking (~1.2s, rect-with-rx interpolation)", {})],
            [("U.S. callout: ", {"bold": True, "color": US_RED}),
             ("'rank N of 20' with leader line", {})],
            "▸ peer tags on Canada, Sweden, U.K., Japan",
            "Takeaway: 'Nine wealthy peers rank above the U.S.'",
        ],
        annotation="Continuity from Step 4 to Step 6 is the whole argument: same world, viewed two ways.",
        talking_point=(
            "This is the closing chart, and the morph is intentional — it "
            "makes visible that the SAME 20 countries we just showed scattered "
            "against inequality are now ranked. The continuity from Step 4 to "
            "Step 6 is the whole argument: same world, viewed two ways."
        ),
        screenshot="step6_league.png",
    ),
    StepSlide(
        step_num=7,
        title="Bonus · A Topography of Opportunity",
        summary="3D extruded county map (deck.gl) — height = mobility for poor-parent children, 1992.",
        bullets=[
            [("Tech: ", {"bold": True, "color": US_RED}),
             ("deck.gl OrbitView + ColumnLayer + PathLayer (state borders)", {})],
            "Lazy-loaded — only downloads when the user scrolls into Step 7",
            "Drag to orbit · Scroll to zoom · Hover for county detail",
            "Mountains and valleys aren't even on the same scale",
        ],
        annotation="2D answers WHICH counties are stuck. 3D answers HOW DRAMATIC the difference is.",
        talking_point=(
            "The 2D choropleth answers which counties are stuck. The 3D "
            "extrusion answers how dramatic the difference is between "
            "neighbors — the mountains and the valleys aren't even on the "
            "same scale. We kept it as a bonus rather than the opener "
            "because color encoding remains more accessible for the "
            "headline argument; the 3D is here to make the magnitude "
            "visceral."
        ),
        screenshot="step7_3d.png",
        bonus=True,
    ),
]


def build_step(prs, step: StepSlide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)

    kicker = (f"BONUS · STEP {step.step_num}" if step.bonus
              else f"STEP {step.step_num}  ·  PART {'I — AMERICAN PROBLEM' if step.step_num <= 3 else 'II — GLOBAL PERSPECTIVE'}")
    slide_header(slide, step.title, kicker=kicker, demo_tag=True)

    # Left column: bullets
    add_text(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.6),
             step.summary, size=14, color=SLATE, font=FONT_TITLE)

    add_bullets(slide, Inches(0.6), Inches(2.85), Inches(5.5), Inches(2.5),
                step.bullets, size=14, color=DARK_INK)

    # Annotation callout
    if step.annotation:
        add_rect(slide, Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.4),
                 RGBColor(0xFD, 0xF2, 0xF1), line_color=US_RED)
        add_rect(slide, Inches(0.6), Inches(5.5), Inches(0.08), Inches(1.4), US_RED)
        add_text(slide, Inches(0.85), Inches(5.65), Inches(5.2), Inches(0.3),
                 "TAKEAWAY", size=10, bold=True, color=US_RED)
        add_text(slide, Inches(0.85), Inches(5.95), Inches(5.2), Inches(0.9),
                 step.annotation, size=13, color=DARK_INK)

    # Right column: screenshot
    add_screenshot_or_placeholder(
        slide, Inches(6.5), Inches(2.0), Inches(6.3), Inches(4.8),
        step.screenshot, f"Step {step.step_num}"
    )

    set_speaker_notes(slide, step.talking_point)


def build_design_decisions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    slide_header(slide, "Design Decisions That Mattered", kicker="WHAT WE LEARNED THE HARD WAY")

    decisions = [
        ("Two-act structure",
         "Local → global lets the reader internalize within-country spread "
         "before weighing the U.S. against peers."),
        ("Sticky right panel",
         "Chart stays anchored; prose advances on the left. Each narrative "
         "beat resolves against a visible mark."),
        ("Word–image binding",
         "When prose names a country, that row is visibly distinguished "
         "(deep-blue label, peer tag, leader line)."),
        ("Curated league of 20",
         "Restricted to major economies. Random low-IGE outliers (Mauritius, "
         "Maldives) muddied the rich-peer framing."),
        ("Button trigger over scroll trigger",
         "Step 0 uses a Play button — fixes layout reflow on fast scrolls."),
        ("3D as bonus, not opener",
         "Color encoding is more accessible for the headline argument. "
         "The 3D is for magnitude impact at the end."),
    ]

    cols = 2
    col_w = Inches(5.95)
    row_h = Inches(1.55)
    start_x = Inches(0.6)
    start_y = Inches(2.0)
    gap_x = Inches(0.3)

    for i, (title, body) in enumerate(decisions):
        col = i % cols
        row = i // cols
        x = start_x + col * (col_w + gap_x)
        y = start_y + row * row_h
        add_rect(slide, x, y, col_w, row_h - Inches(0.15),
                 RGBColor(0xFF, 0xFF, 0xFF), line_color=BORDER_GRAY)
        add_rect(slide, x, y, Inches(0.08), row_h - Inches(0.15), US_RED)
        add_text(slide, x + Inches(0.3), y + Inches(0.18), col_w - Inches(0.4),
                 Inches(0.4), f"{i+1}. {title}", size=14, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.3), y + Inches(0.6), col_w - Inches(0.4),
                 row_h - Inches(0.7), body, size=11, color=DARK_INK)

    set_speaker_notes(
        slide,
        "Most of these decisions came from things that broke. The button "
        "trigger replaced a scroll trigger that raced. The curated league "
        "replaced a naive 'all countries' view that buried the U.S. in noise. "
        "The 3D was deliberately kept as a bonus rather than competing with "
        "the 2D choropleth's accessibility for the headline argument.",
    )


def build_animation_highlights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    slide_header(slide, "Animation & Interaction Highlights",
                 kicker="WHAT MOVES, WHAT THE READER CAN DO")

    highlights = [
        ("Morph animation", "scatter → bar ranking on Step 6"),
        ("Animated reveal", "county-by-county fill transition on Step 2"),
        ("Button-triggered playback", "Step 0 with Replay"),
        ("Hover tooltips", "every chart, with country / county / value detail"),
        ("State drill-down", "in the bubble (Dorling) view"),
        ("3D orbit + zoom", "bonus Step 7 (deck.gl OrbitView)"),
        ("In-chart annotations", "every step — prose references resolve to specific marks"),
    ]

    top = Inches(2.0)
    row_h = Inches(0.65)
    for i, (head, body) in enumerate(highlights):
        y = top + i * row_h
        add_rect(slide, Inches(0.6), y, Inches(0.45), Inches(0.45), US_RED)
        add_text(slide, Inches(0.65), y + Inches(0.05), Inches(0.4),
                 Inches(0.4), str(i + 1), size=14, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1.25), y + Inches(0.02), Inches(3.8),
                 Inches(0.5), head, size=14, bold=True, color=NAVY)
        add_text(slide, Inches(5.0), y + Inches(0.05), Inches(7.5),
                 Inches(0.5), body, size=13, color=DARK_INK)

    set_speaker_notes(
        slide,
        "Every chart has at least one motion or interaction. Step 6's morph "
        "is the visual climax — the same 20 countries from Step 4's scatter "
        "literally flow into a sorted ranking. Step 7's 3D adds a magnitude "
        "dimension that flat color can't carry.",
    )


def build_reflection(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_BG)
    slide_header(slide, "Reflection", kicker="WHAT WORKED · WHAT WE'D REDO")

    cards = [
        (
            "What worked",
            NORDIC_GREEN,
            "Scrollytelling kept the argument sequenced. Sticky charts "
            "gave each prose beat a visual anchor. Adding the deck.gl 3D "
            "bonus at the end gave a 'magnitude reveal' without "
            "competing with the 2D choropleth's accessibility for the "
            "headline argument.",
        ),
        (
            "What we'd redo",
            DATA_BLUE,
            "Invest more in keyboard-only interaction for the 3D scene "
            "— orbit/zoom currently rely on mouse gestures. Add an "
            "elevation slider for non-mouse users.",
        ),
        (
            "Hardest part",
            US_RED,
            "Keeping the same 20-country roster visually consistent "
            "across Steps 4, 5, 6. The morph from scatter to ranking is "
            "what makes that roster legible — the eye carries identity "
            "from one chart to the next.",
        ),
    ]

    card_w = Inches(3.95)
    card_h = Inches(4.0)
    start_x = Inches(0.6)
    top = Inches(2.0)
    gap = Inches(0.2)
    for i, (title, color, body) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, top, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF),
                 line_color=BORDER_GRAY)
        add_rect(slide, x, top, card_w, Inches(0.5), color)
        add_text(slide, x + Inches(0.3), top + Inches(0.1), card_w - Inches(0.4),
                 Inches(0.35), title, size=15, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
        add_text(slide, x + Inches(0.3), top + Inches(0.85),
                 card_w - Inches(0.4), card_h - Inches(1.0),
                 body, size=13, color=DARK_INK)

    set_speaker_notes(
        slide,
        "If we had another week we'd add keyboard navigation to the 3D "
        "scene and an elevation slider for non-mouse users. The hardest "
        "engineering problem was actually keeping the 20-country roster "
        "consistent across Steps 4, 5, and 6 — the morph animation is "
        "what carries identity through.",
    )


def make_qr(url, fname):
    img = qrcode.make(url)
    out = QR_DIR / fname
    img.save(out)
    return out


def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)

    add_text(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.5),
             "Q & A", size=22, bold=True,
             color=RGBColor(0xFF, 0xC0, 0xB6), font=FONT_TITLE,
             align=PP_ALIGN.LEFT)

    add_text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.7),
             "Thank you.", size=44, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), font=FONT_TITLE)

    add_text(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(0.5),
             "We're happy to take questions.", size=18,
             color=RGBColor(0xE0, 0xE6, 0xEE), font=FONT_TITLE)

    # QR + URLs
    qr_size = Inches(2.0)
    qr_y = Inches(4.0)

    qr_live = make_qr(LIVE_URL, "live.png")
    qr_repo = make_qr(REPO_URL, "repo.png")

    slide.shapes.add_picture(str(qr_live), Inches(1.5), qr_y, qr_size, qr_size)
    add_text(slide, Inches(1.5), qr_y + qr_size + Inches(0.1), Inches(2.0),
             Inches(0.3), "LIVE DEMO", size=11, bold=True,
             color=RGBColor(0xFF, 0xC0, 0xB6), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.0), qr_y + qr_size + Inches(0.4), Inches(3.0),
             Inches(0.3), "57777-bit.github.io/5609finalProject",
             size=10, color=RGBColor(0xE0, 0xE6, 0xEE), font=FONT_MONO,
             align=PP_ALIGN.CENTER)

    slide.shapes.add_picture(str(qr_repo), Inches(8.5), qr_y, qr_size, qr_size)
    add_text(slide, Inches(8.5), qr_y + qr_size + Inches(0.1), Inches(2.0),
             Inches(0.3), "REPO", size=11, bold=True,
             color=RGBColor(0xFF, 0xC0, 0xB6), align=PP_ALIGN.CENTER)
    add_text(slide, Inches(8.0), qr_y + qr_size + Inches(0.4), Inches(3.0),
             Inches(0.3), "github.com/57777-bit/5609finalProject",
             size=10, color=RGBColor(0xE0, 0xE6, 0xEE), font=FONT_MONO,
             align=PP_ALIGN.CENTER)

    # Acknowledgements
    add_text(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.3),
             "ACKNOWLEDGEMENTS", size=10, bold=True,
             color=RGBColor(0xFF, 0xC0, 0xB6))
    add_text(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
             "Chetty et al. · Opportunity Atlas   ·   World Bank GDIM 2023   ·   "
             "Frederick Solt SWIID v9.6   ·   Mike Bostock TopoJSON Atlas   ·   D3   ·   Svelte/SvelteKit",
             size=10, color=RGBColor(0xE0, 0xE6, 0xEE), font=FONT_BODY)

    set_speaker_notes(
        slide,
        "Open the floor. The QR codes go to the live demo and the repo. "
        "Common questions to be ready for: data limitations, why we picked "
        "1992 specifically (Opportunity Atlas's earliest robust cohort), "
        "why IGE not rank correlation (GDIM uses IGE; we kept their measure "
        "to avoid double conversion), and why the 3D is a bonus not the "
        "opener (color encoding is more accessible for the headline argument).",
    )


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title(prs)
    build_question(prs)
    build_data_sources(prs)
    build_tech_stack(prs)
    for step in STEPS:
        build_step(prs, step)
    build_design_decisions(prs)
    build_animation_highlights(prs)
    build_reflection(prs)
    build_closing(prs)

    # Add page numbers + footer to all slides except the title and closing
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == total - 1:
            continue
        slide_footer(slide, i + 1, total)

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}  ({total} slides)")


if __name__ == "__main__":
    main()
